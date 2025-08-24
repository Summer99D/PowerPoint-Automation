from pptx import Presentation
# *************************
# Textbox lister 
# **************************

# def list_text_boxes(presentation, slide_num):
#    slide = presentation.slides[slide_num - 1]  # slides are 0-indexed
#    text_boxes = []
#    for shape in slide.shapes:
#        if shape.has_text_frame:
#            if shape.text.strip():  # avoid empty text boxes
#                text_boxes.append(shape.text.strip())
#   return text_boxes

#if __name__ == "__main__":
    # Change this to your PPTX file
#    pptx_file = "powerpoint/template.pptx"

    # Load the presentation
#    prs = Presentation(pptx_file)

#    # Example: print all text from slide X
#    slide_num = 7
#    texts = list_text_boxes(prs, slide_num)

#    print(f"Slide {slide_num} has {len(texts)} text box(es):")
#    for i, txt in enumerate(texts, 1):
#        print(f"{i}: {txt}")


# *************************
# shape lister: this is a better function compared to text box one since it shows both shapes ad texts!
# **************************

def list_shapes(pptx_file, slide_num):
    prs = Presentation(pptx_file)
    slide = prs.slides[slide_num - 1]
    
    for idx, shape in enumerate(slide.shapes, start=1):
        shape_type = shape.shape_type
        shape_name = shape.name
        has_text = shape.has_text_frame
        text = shape.text if has_text and shape.text else ""
        
        print(f"{idx}. Name: {shape_name}, Type: {shape_type}, Text: '{text}'")

# Example usage:
list_shapes("powerpoint/template.pptx", 3)