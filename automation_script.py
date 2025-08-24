import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import argparse
import os
from datetime import datetime
import seaborn as sns
import numpy as np

# -------------------
# Helper Functions
# -------------------
def validate_date(date_str):
    """Validate date format and ensure it's within dataset range (2014-12-22 to 2024-12-20)."""
    try:
        date = pd.to_datetime(date_str)
        if date < pd.to_datetime("2014-12-22") or date > pd.to_datetime("2024-12-20"):
            raise ValueError("Date must be between 2014-12-22 and 2024-12-20.")
        return date
    except ValueError as e:
        raise ValueError(f"Invalid date format or range: {e}")

def update_text_of_textbox(presentation, slide_num, shape_name, new_text):
    """Update the text of a specific textbox or shape with text while preserving its formatting."""
    slide = presentation.slides[slide_num - 1]
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            text_frame = shape.text_frame
            first_paragraph = text_frame.paragraphs[0]
            first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
            # Preserve formatting of the first run
            font = first_run.font
            font_name = font.name
            font_size = font.size
            font_bold = font.bold
            font_italic = font.italic
            font_underline = font.underline
            # Handle undefined color gracefully
            font_color = font.color.rgb if font.color.type is not None and font.color.rgb else RGBColor(0, 0, 0)  # Default to black
            # Clear existing text and apply new text with preserved formatting
            text_frame.clear()
            new_run = text_frame.paragraphs[0].add_run()
            new_run.text = new_text
            # Reapply formatting
            new_run.font.name = font_name
            new_run.font.size = font_size
            new_run.font.bold = font_bold
            new_run.font.italic = font_italic
            new_run.font.underline = font_underline
            new_run.font.color.rgb = font_color
            return
    print(f"Warning: Shape {shape_name} not found on slide {slide_num}.")

def get_placeholder_position(presentation, slide_num, shape_name):
    """Get the position (left, top, width, height) of a shape by name."""
    slide = presentation.slides[slide_num - 1]
    for shape in slide.shapes:
        if shape.name == shape_name and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return shape.left, shape.top, shape.width, shape.height
    print(f"Warning: Shape {shape_name} not found on slide {slide_num}.")
    return Inches(0.5), Inches(1), Inches(9), Inches(4.5)  # Fallback position

# -------------------
# Slide 3: Performance Comparison
# -------------------
def compare_with_sp500(slide, date, symbol, merged_df, presentation):
    """Add text comparing stock's adjusted closing price to average S&P 500 stock price."""
    target_date = pd.to_datetime(date)
    date_data = merged_df[merged_df['Date'] == target_date]
    used_date = target_date
    if date_data.empty:
        print(f"No data available for {date}, using closest date.")
        available_dates = merged_df['Date'].unique()
        closest_date = min(available_dates, key=lambda d: abs(pd.to_datetime(d) - target_date))
        date_data = merged_df[merged_df['Date'] == closest_date]
        used_date = closest_date
        print(f"Using closest date: {used_date.strftime('%Y-%m-%d')}")
    
    stock_price = date_data[date_data['Symbol'] == symbol]['Adj Close'].iloc[0] if not date_data[date_data['Symbol'] == symbol].empty else None
    avg_stock_price = date_data['Adj Close'].mean()

    # Check for NaN or missing data
    if pd.isna(stock_price) or pd.isna(avg_stock_price) or stock_price is None:
        print(f"Warning: Invalid data - Stock Price: {stock_price}, S&P 500 Average: {avg_stock_price}")
        text = f"Data unavailable for {symbol} on {used_date.strftime('%Y-%m-%d')} (invalid values)"
    else:
        performance = "outperformed" if stock_price > avg_stock_price else "underperformed" if stock_price < avg_stock_price else "matched"
        text = f"The price of {symbol} on {used_date.strftime('%Y-%m-%d')} is ${stock_price:.2f}. The price {performance} S&P 500 average."

    # Update textboxes
    update_text_of_textbox(presentation, 3, "TextBox 1", text)  # Comparison sample text
    return used_date  # Return the actual date used

# -------------------
# Slide 4: Scatterplot Comparing Stock to Others
# -------------------
def add_scatter_chart(slide, date, symbol, merged_df, presentation):
    """Add a scatterplot comparing the stock's adjusted closing price to all other stocks in Rectangle 1, with outlier labels."""
    target_date = pd.to_datetime(date)
    date_data = merged_df[merged_df['Date'] == target_date]
    used_date = target_date
    if date_data.empty:
        print(f"No data available for {date}, using closest date.")
        available_dates = merged_df['Date'].unique()
        closest_date = min(available_dates, key=lambda d: abs(pd.to_datetime(d) - target_date))
        date_data = merged_df[merged_df['Date'] == closest_date]
        used_date = closest_date
        print(f"Using closest date: {used_date.strftime('%Y-%m-%d')}")

    # Check for NaN in Adj Close column
    if date_data['Adj Close'].isna().any():
        print(f"Warning: {date_data['Adj Close'].isna().sum()} NaN values in Adj Close column for {used_date.strftime('%Y-%m-%d')}")

    # Identify outliers (±2 standard deviations)
    prices = date_data['Adj Close']
    mean_price = prices.mean()
    std_price = prices.std()
    outlier_threshold = 2  # Standard deviations
    outliers = date_data[(prices > mean_price + outlier_threshold * std_price) | (prices < mean_price - outlier_threshold * std_price)]

    # Create scatterplot
    plt.figure(figsize=(14, 7))
    x = range(len(date_data))  # Index for each stock
    colors = ['#ff7f0e' if s == symbol else '#1f77b4' for s in date_data['Symbol']]
    alphas = [1.0 if s == symbol else 0.5 for s in date_data['Symbol']]  # 100% opacity for selected stock, 50% for others
    for i, (xi, yi, color, alpha) in enumerate(zip(x, date_data['Adj Close'], colors, alphas)):
        plt.scatter(xi, yi, c=color, alpha=alpha, s=50)
        # Add labels for outliers
        if date_data.iloc[i]['Symbol'] in outliers['Symbol'].values:
            plt.text(xi, yi, date_data.iloc[i]['Symbol'], fontsize=8, ha='center', va='bottom')
    plt.title(f"Adjusted Closing Prices of S&P 500 Stocks on {used_date.strftime('%Y-%m-%d')}\n{symbol} Highlighted", fontsize=12)
    plt.xlabel('Stocks', fontsize=10)
    plt.ylabel('Adjusted Closing Price (USD)', fontsize=10)
    plt.xticks([]) 
    plt.grid(True, axis='y')
    plt.legend(loc='upper left')  # adding legend
    plt.tight_layout()

    os.makedirs('scratch', exist_ok=True)
    chart_path = f'scratch/{symbol}_scatter_chart.png'
    plt.savefig(chart_path, dpi=100)
    plt.close()

    # Insert chart at Rectangle 1 position
    left, top, width, height = get_placeholder_position(presentation, 4, "Rectangle 1")
    slide.shapes.add_picture(chart_path, left, top, width, height)
    return used_date  # Return the actual date used

# -------------------
# Slide 5: Weekly Chart
# -------------------
def add_weekly_chart(slide, date, symbol, merged_df, presentation):
    """Add two line charts: stock in Rectangle 1, S&P 500 in Rectangle 5."""
    target_date = pd.to_datetime(date)
    start_date = target_date - pd.Timedelta(days=7)
    data = merged_df[(merged_df['Date'] >= start_date) & (merged_df['Date'] <= target_date) & ((merged_df['Symbol'] == symbol) | (merged_df['Symbol'] == merged_df['Symbol'].iloc[0]))][['Date', 'Symbol', 'Adj Close', 'S&P500']].sort_values('Date')

    if data[data['Symbol'] == symbol].empty:
        print(f"Warning: No data for {symbol} in weekly period ({start_date.strftime('%Y-%m-%d')} to {target_date.strftime('%Y-%m-%d')}). Skipping chart.")
        return

    # Check for NaN in data
    if data['Adj Close'].isna().any() or data['S&P500'].isna().any():
        print(f"Warning: NaN values in weekly data - Stock: {data[data['Symbol'] == symbol]['Adj Close'].isna().sum()}, S&P 500: {data['S&P500'].isna().sum()}")

    # Stock chart (Rectangle 1)
    plt.figure(figsize=(10, 6))
    plt.plot(data[data['Symbol'] == symbol]['Date'], data[data['Symbol'] == symbol]['Adj Close'], label=f'{symbol} Adj Close', color='#1f77b4', linewidth=2)
    plt.title(f'{symbol} Weekly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Adjusted Closing Price (USD)', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_stock = f'scratch/{symbol}_weekly_stock.png'
    plt.savefig(chart_path_stock, dpi=100)
    plt.close()

    # S&P 500 chart (Rectangle 5)
    plt.figure(figsize=(10, 6))
    plt.plot(data['Date'], data['S&P500'], label='S&P 500', color='#ff7f0e', linewidth=2)
    plt.title('S&P 500 Weekly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Index Value', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_sp500 = f'scratch/{symbol}_weekly_sp500.png'
    plt.savefig(chart_path_sp500, dpi=100)
    plt.close()

    # Insert charts
    left, top, width, height = get_placeholder_position(presentation, 5, "Rectangle 1")
    slide.shapes.add_picture(chart_path_stock, left, top, width, height)
    left, top, width, height = get_placeholder_position(presentation, 5, "Rectangle 5")
    slide.shapes.add_picture(chart_path_sp500, left, top, width, height)

# -------------------
# Slide 6: Monthly Chart
# -------------------
def add_monthly_chart(slide, date, symbol, merged_df, presentation):
    """Add two line charts: stock in Rectangle 1, S&P 500 in Rectangle 5."""
    target_date = pd.to_datetime(date)
    start_date = target_date - pd.Timedelta(days=30)
    data = merged_df[(merged_df['Date'] >= start_date) & (merged_df['Date'] <= target_date) & ((merged_df['Symbol'] == symbol) | (merged_df['Symbol'] == merged_df['Symbol'].iloc[0]))][['Date', 'Symbol', 'Adj Close', 'S&P500']].sort_values('Date')

    if data[data['Symbol'] == symbol].empty:
        print(f"Warning: No data for {symbol} in monthly period ({start_date.strftime('%Y-%m-%d')} to {target_date.strftime('%Y-%m-%d')}). Skipping chart.")
        return

    # Check for NaN in data
    if data['Adj Close'].isna().any() or data['S&P500'].isna().any():
        print(f"Warning: NaN values in monthly data - Stock: {data[data['Symbol'] == symbol]['Adj Close'].isna().sum()}, S&P 500: {data['S&P500'].isna().sum()}")

    # Stock chart (Rectangle 1)
    plt.figure(figsize=(10, 6))
    plt.plot(data[data['Symbol'] == symbol]['Date'], data[data['Symbol'] == symbol]['Adj Close'], label=f'{symbol} Adj Close', color='#1f77b4', linewidth=2)
    plt.title(f'{symbol} Monthly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Adjusted Closing Price (USD)', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_stock = f'scratch/{symbol}_monthly_stock.png'
    plt.savefig(chart_path_stock, dpi=100)
    plt.close()

    # S&P 500 chart (Rectangle 5)
    plt.figure(figsize=(10, 6))
    plt.plot(data['Date'], data['S&P500'], label='S&P 500', color='#ff7f0e', linewidth=2)
    plt.title('S&P 500 Monthly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Index Value', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_sp500 = f'scratch/{symbol}_monthly_sp500.png'
    plt.savefig(chart_path_sp500, dpi=100)
    plt.close()

    # Insert charts
    left, top, width, height = get_placeholder_position(presentation, 6, "Rectangle 1")
    slide.shapes.add_picture(chart_path_stock, left, top, width, height)
    left, top, width, height = get_placeholder_position(presentation, 6, "Rectangle 5")
    slide.shapes.add_picture(chart_path_sp500, left, top, width, height)

# -------------------
# Slide 7: Yearly Chart
# -------------------
def add_yearly_chart(slide, date, symbol, merged_df, presentation):
    """Add two line charts: stock in Rectangle 1, S&P 500 in Rectangle 5."""
    target_date = pd.to_datetime(date)
    start_date = target_date - pd.Timedelta(days=365)
    data = merged_df[(merged_df['Date'] >= start_date) & (merged_df['Date'] <= target_date) & ((merged_df['Symbol'] == symbol) | (merged_df['Symbol'] == merged_df['Symbol'].iloc[0]))][['Date', 'Symbol', 'Adj Close', 'S&P500']].sort_values('Date')

    if data[data['Symbol'] == symbol].empty:
        print(f"Warning: No data for {symbol} in yearly period ({start_date.strftime('%Y-%m-%d')} to {target_date.strftime('%Y-%m-%d')}). Skipping chart.")
        return

    # Check for NaN in data
    if data['Adj Close'].isna().any() or data['S&P500'].isna().any():
        print(f"Warning: NaN values in yearly data - Stock: {data[data['Symbol'] == symbol]['Adj Close'].isna().sum()}, S&P 500: {data['S&P500'].isna().sum()}")

    # Stock chart (Rectangle 1)
    plt.figure(figsize=(10, 6))
    plt.plot(data[data['Symbol'] == symbol]['Date'], data[data['Symbol'] == symbol]['Adj Close'], label=f'{symbol} Adj Close', color='#1f77b4', linewidth=2)
    plt.title(f'{symbol} Yearly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Adjusted Closing Price (USD)', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_stock = f'scratch/{symbol}_yearly_stock.png'
    plt.savefig(chart_path_stock, dpi=100)
    plt.close()

    # S&P 500 chart (Rectangle 5)
    plt.figure(figsize=(10, 6))
    plt.plot(data['Date'], data['S&P500'], label='S&P 500', color='#ff7f0e', linewidth=2)
    plt.title('S&P 500 Yearly Performance', fontsize=12)
    plt.xlabel('Date', fontsize=10)
    plt.ylabel('Index Value', fontsize=10)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    chart_path_sp500 = f'scratch/{symbol}_yearly_sp500.png'
    plt.savefig(chart_path_sp500, dpi=100)
    plt.close()

    # Insert charts
    left, top, width, height = get_placeholder_position(presentation, 7, "Rectangle 1")
    slide.shapes.add_picture(chart_path_stock, left, top, width, height)
    left, top, width, height = get_placeholder_position(presentation, 7, "Rectangle 5")
    slide.shapes.add_picture(chart_path_sp500, left, top, width, height)

# -------------------
# Main Function
# -------------------
def main(date, symbol):
    """Main function to generate the presentation."""
    # Load presentation template
    presentation = Presentation(os.path.join("powerpoint", "template.pptx"))

    # Diagnostic: Print shape names and text for Slides 5-7 before updates
    for slide_num in [5, 6, 7]:
        print(f"\nSlide {slide_num} shapes before updates:")
        slide = presentation.slides[slide_num - 1]
        for shape in slide.shapes:
            shape_type = 'AUTO_SHAPE' if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE else 'TEXT_BOX'
            text = shape.text if shape.has_text_frame else ''
            print(f"  Name: {shape.name}, Type: {shape_type}, Text: '{text}'")

    # Load dataset
    try:
        merged_df = pd.read_csv('datasets/merged_sp500.csv')
        merged_df['Date'] = pd.to_datetime(merged_df['Date'])
    except FileNotFoundError:
        print("Error: merged_sp500.csv not found in datasets directory.")
        return

    # Validate inputs
    try:
        target_date = validate_date(date)
        if symbol not in merged_df['Symbol'].unique():
            print(f"Error: Symbol {symbol} not found in dataset.")
            return
    except ValueError as e:
        print(f"Error: {e}")
        return

    # Update slides and get the actual date used
    used_date = compare_with_sp500(presentation.slides[2], target_date.strftime("%Y-%m-%d"), symbol, merged_df, presentation)
    used_date = add_scatter_chart(presentation.slides[3], used_date.strftime("%Y-%m-%d") if used_date != target_date else date, symbol, merged_df, presentation)
    add_weekly_chart(presentation.slides[4], used_date.strftime("%Y-%m-%d") if used_date != target_date else date, symbol, merged_df, presentation)
    add_monthly_chart(presentation.slides[5], used_date.strftime("%Y-%m-%d") if used_date != target_date else date, symbol, merged_df, presentation)
    add_yearly_chart(presentation.slides[6], used_date.strftime("%Y-%m-%d") if used_date != target_date else date, symbol, merged_df, presentation)

    # Diagnostic: Print shape names and text for Slides 5-7 after updates
    for slide_num in [5, 6, 7]:
        print(f"\nSlide {slide_num} shapes after updates:")
        slide = presentation.slides[slide_num - 1]
        for shape in slide.shapes:
            shape_type = 'AUTO_SHAPE' if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE else 'TEXT_BOX'
            text = shape.text if shape.has_text_frame else ''
            print(f"  Name: {shape.name}, Type: {shape_type}, Text: '{text}'")

    # Save the presentation with the actual date used
    output_path = os.path.join("powerpoint", f"{symbol}_{used_date.strftime('%Y-%m-%d')}_SUMMARY.pptx")
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate stock performance presentation.")
    parser.add_argument("--date", required=True, help="Date in YYYY-MM-DD format")
    parser.add_argument("--symbol", required=True, help="Stock symbol (e.g., AAPL)")
    args = parser.parse_args()
    main(args.date, args.symbol)